import os
import io
import streamlit as st
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches

def convert_docx_to_pptx(docx_file):
    doc = Document(docx_file)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]   # Layout trang bìa
    content_slide_layout = prs.slide_layouts[1] # Layout nội dung

    current_slide = None
    text_box = None
    title_slide = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Header
        if text.lower().startswith("header:"):
            title_text = text.split(":", 1)[1].strip()
            title_slide = prs.slides.add_slide(title_slide_layout)
            title = title_slide.shapes.title
            title.text = title_text
            for run in title.text_frame.paragraphs[0].runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(36)

        # Sub
        elif text.lower().startswith("sub:") and title_slide:
            sub_text = text.split(":", 1)[1].strip()
            subtitle = title_slide.placeholders[1]
            subtitle.text = sub_text
            for run in subtitle.text_frame.paragraphs[0].runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(24)

        # Info
        elif text.lower().startswith("info:") and title_slide:
            info_text = text.split(":", 1)[1].strip()
            left = Inches(1)
            top = Inches(5.5)
            width = Inches(8)
            height = Inches(0.5)
            textbox = title_slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = info_text
            for run in tf.paragraphs[0].runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(14)
                run.font.italic = True

        # Slide nội dung
        elif text.lower().startswith("slide "):
            current_slide = prs.slides.add_slide(content_slide_layout)
            if ":" in text:
                slide_title = text.split(":", 1)[1].strip()
            else:
                slide_title = text
            current_slide.shapes.title.text = slide_title
            for run in current_slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(28)

            text_box = current_slide.placeholders[1].text_frame
            text_box.clear()

        # Bullet
        else:
            if current_slide and text_box:
                p = text_box.add_paragraph()
                p.text = text.lstrip("•- ")
                p.level = 0
                for run in p.runs:
                    run.font.name = "Segoe UI"
                    run.font.size = Pt(20)

    # Lưu vào bộ nhớ
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ================================
# Streamlit UI
# ================================
st.set_page_config(page_title="Tins_SlideForge_Pro2", page_icon="📊")

st.title("📄 ➜ 📊 Word to PowerPoint Converter")
st.write("Upload a `.docx` file and convert it into a PowerPoint presentation.")

uploaded_file = st.file_uploader("Upload Word (.docx)", type=["docx"])

if uploaded_file is not None:
    if st.button("Convert to PowerPoint"):
        try:
            pptx_file = convert_docx_to_pptx(uploaded_file)
            st.success("✅ Conversion successful!")

            st.download_button(
                label="📥 Download PowerPoint",
                data=pptx_file,
                file_name=os.path.splitext(uploaded_file.name)[0] + ".pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        except Exception as e:
            st.error(f"❌ Error: {str(e)}")

with st.expander("📌 Hướng dẫn"):
    st.markdown("""
    **QUY TẮC NHẬP FILE WORD (.docx):**

    1. **Trang bìa:**
       - `Header:` tiêu đề chính  
       - `Sub:` phụ đề (tùy chọn)  
       - `Info:` thông tin thêm (người trình bày, ngày tháng...)  

    2. **Slide nội dung:**
       - Bắt đầu bằng `Slide X: ...`  
       - Ví dụ: `Slide 1: Giới thiệu`  

    3. **Bullet point:**
       - Dùng dấu `-` hoặc `•` ở đầu dòng  
       - Mỗi bullet 1 dòng riêng  

    👉 Toàn bộ slide sẽ dùng font **Segoe UI**.
    """)
